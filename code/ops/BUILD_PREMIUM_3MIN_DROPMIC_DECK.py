from __future__ import annotations

import json
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[2]
OUT_OPS = ROOT / "out" / "ops"
MISSION_DIR = OUT_OPS / "investor_mission_control"
MASTER_VAL_PATH = OUT_OPS / "master_valuation" / "master_valuation_latest.json"
LIC_SCENARIOS_PATH = OUT_OPS / "master_valuation" / "valuation_licensing_scenarios_latest.json"
MISSION_PACK_PATH = MISSION_DIR / "investor_mission_control_pack_latest.json"
PANEL_PATH = OUT_OPS / "live_breadth_value_panel_latest.json"

Presentation = None
RGBColor = None
MSO_AUTO_SHAPE_TYPE = None
PP_ALIGN = None
Inches = None
Pt = None


def _load_pptx_imports() -> tuple[bool, str]:
    global Presentation, RGBColor, MSO_AUTO_SHAPE_TYPE, PP_ALIGN, Inches, Pt
    try:
        from pptx import Presentation as _Presentation
        from pptx.dml.color import RGBColor as _RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN
        from pptx.util import Inches as _Inches
        from pptx.util import Pt as _Pt

        Presentation = _Presentation
        RGBColor = _RGBColor
        MSO_AUTO_SHAPE_TYPE = _MSO_AUTO_SHAPE_TYPE
        PP_ALIGN = _PP_ALIGN
        Inches = _Inches
        Pt = _Pt
        return True, "ok"
    except Exception as exc:
        return False, str(exc)


def ensure_pptx_runtime() -> dict[str, Any]:
    ok, reason = _load_pptx_imports()
    if ok:
        return {
            "available": True,
            "source": "installed",
            "reason": "ok",
        }

    install_cmd = [sys.executable, "-m", "pip", "install", "python-pptx"]
    try:
        subprocess.run(install_cmd, check=True, capture_output=True, text=True)
    except Exception as install_exc:
        return {
            "available": False,
            "source": "missing",
            "reason": f"import_error={reason}; install_error={install_exc}",
        }

    ok2, reason2 = _load_pptx_imports()
    if ok2:
        return {
            "available": True,
            "source": "auto_installed",
            "reason": "ok",
        }

    return {
        "available": False,
        "source": "missing",
        "reason": f"post_install_import_error={reason2}",
    }


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def utc_tag() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")


def load_json(path: Path) -> Any:
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None


def safe_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except Exception:
        return default


def safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(float(value))
    except Exception:
        return default


def money(value: Any) -> str:
    return f"${safe_float(value):,.2f}"


def short_money(value: Any) -> str:
    amount = safe_float(value)
    if abs(amount) >= 1_000_000_000:
        return f"${amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"${amount / 1_000_000:.2f}M"
    if abs(amount) >= 1_000:
        return f"${amount / 1_000:.2f}K"
    return money(amount)


def latest_patent_anchor() -> dict[str, str]:
    ops_parent = ROOT.parent / "out" / "ops"
    if not ops_parent.exists():
        return {}
    candidates = sorted(ops_parent.glob("patent_filing_tracker_*.json"))
    if not candidates:
        return {}
    payload = load_json(candidates[-1])
    if not isinstance(payload, dict):
        return {}
    anchor = payload.get("application_anchor", {})
    return anchor if isinstance(anchor, dict) else {}


def build_payload() -> dict[str, Any]:
    master = load_json(MASTER_VAL_PATH)
    scenarios = load_json(LIC_SCENARIOS_PATH)
    mission = load_json(MISSION_PACK_PATH)
    panel = load_json(PANEL_PATH)

    if not isinstance(master, dict):
        raise RuntimeError(f"Missing master valuation payload: {MASTER_VAL_PATH}")
    if not isinstance(scenarios, dict):
        raise RuntimeError(f"Missing licensing scenarios payload: {LIC_SCENARIOS_PATH}")

    valuation = master.get("valuation", {}) if isinstance(master, dict) else {}
    inputs = master.get("inputs", {}) if isinstance(master, dict) else {}
    headline = (panel.get("headline", {}) or {}) if isinstance(panel, dict) else {}
    mission_headline = (mission.get("headline", {}) or {}) if isinstance(mission, dict) else {}
    pitch_segments = (
        ((mission.get("three_min_nobel_pitch", {}) or {}).get("segments", []))
        if isinstance(mission, dict)
        else []
    )
    patent_anchor = latest_patent_anchor()

    annual_value = max(
        safe_float(inputs.get("annual_value_signal_usd"), 0.0),
        safe_float(headline.get("total_estimated_annual_value_usd"), 0.0),
        safe_float(mission_headline.get("annual_value_signal_usd"), 0.0),
    )
    measured_sources = max(
        safe_int(inputs.get("measured_sources"), 0),
        safe_int(headline.get("measured_sources"), 0),
        safe_int(mission_headline.get("measured_sources"), 0),
    )
    enabled_sources = max(
        safe_int(inputs.get("enabled_sources"), 0),
        safe_int(headline.get("enabled_sources"), 0),
        safe_int(mission_headline.get("enabled_sources"), 0),
    )
    coverage_pct = max(
        safe_float(inputs.get("measured_coverage_pct"), 0.0),
        safe_float(headline.get("measured_coverage_pct"), 0.0),
        safe_float(mission_headline.get("measured_coverage_pct"), 0.0),
    )

    benchmark_gain_pct = safe_float(headline.get("cross_sector_recommended_prevented_pct"), 0.0)
    router_edge_pct = safe_float(inputs.get("router_edge_pct"), safe_float(mission_headline.get("router_edge_pct"), 0.0))
    harmonic_win_rate_pct = safe_float(inputs.get("harmonic_win_rate_pct"), safe_float(mission_headline.get("harmonic_win_rate_pct"), 0.0))

    top_sector = str(headline.get("top_sector") or mission_headline.get("top_sector") or "financial_market_infra")
    top_sector_hourly = max(
        safe_float(headline.get("top_sector_hourly_value_usd"), 0.0),
        safe_float(mission_headline.get("top_sector_hourly_value_usd"), 0.0),
    )

    valuation_proxy = safe_float(valuation.get("master_valuation_proxy_usd"), 0.0)
    valuation_increment = safe_float(valuation.get("valuation_increment_usd"), 0.0)
    grant_license = safe_float(valuation.get("grant_finding_and_ranking_system_license_value_usd"), 0.0)
    grant_pipeline = safe_float(valuation.get("grant_and_opportunity_pipeline_value_usd"), 0.0)
    digital_scout = safe_float(valuation.get("digital_scout_value_usd"), 0.0)
    trading_value = safe_float(valuation.get("institutional_trading_system_value_usd"), 0.0)
    autonomy_value = safe_float(valuation.get("validated_engine_autonomy_value_usd"), 0.0)

    scenario_label = str(scenarios.get("recommended_scenario_label") or "")
    scenario_arr = safe_float(scenarios.get("recommended_arr_usd"), 0.0)
    scenario_year1 = safe_float(scenarios.get("recommended_year1_revenue_usd"), 0.0)
    scenario_ev_low = safe_float(scenarios.get("recommended_implied_ev_low_usd"), 0.0)
    scenario_ev_high = safe_float(scenarios.get("recommended_implied_ev_high_usd"), 0.0)

    point_01pct_value = annual_value * 0.0001
    twenty_lane_billion_surface = 20_000_000_000.0
    point_01pct_twenty_lanes = twenty_lane_billion_surface * 0.0001

    claim_lines = [
        f"Measured live source coverage: {measured_sources}/{enabled_sources} ({coverage_pct:.2f}%).",
        f"Measured annual preserved-value surface: {short_money(annual_value)}.",
        f"Cross-sector modeled prevention benchmark: {benchmark_gain_pct:.2f}%.",
        f"Router edge: {router_edge_pct:.2f}% | harmonic win-rate: {harmonic_win_rate_pct:.2f}%.",
        f"0.01% of current measured annual surface equals {short_money(point_01pct_value)} per year.",
        f"If 20 lanes each reached $1B annual surface, 0.01% equals {short_money(point_01pct_twenty_lanes)} per year.",
    ]

    segments = [
        {
            "start": "00:00",
            "end": "00:30",
            "title": "The Drop-Mic Opening",
            "bullets": [
                "Founder of the year energy: this is not a concept deck, this is a live operating stack.",
                f"We measure {measured_sources}/{enabled_sources} live sources with {coverage_pct:.2f}% coverage.",
                f"Today the measured annual value surface is {short_money(annual_value)}.",
            ],
            "script": (
                "I am not pitching an idea. I am showing a measured operating system. "
                f"Right now we measure {measured_sources} live sources out of {enabled_sources}, "
                f"with {coverage_pct:.2f}% live coverage and {short_money(annual_value)} annual preserved-value surface."
            ),
        },
        {
            "start": "00:30",
            "end": "01:00",
            "title": "Proof and Benchmark",
            "bullets": [
                f"Cross-sector modeled prevention benchmark: {benchmark_gain_pct:.2f}%.",
                f"Top lane: {top_sector} at {short_money(top_sector_hourly)}/hour equivalent value pressure.",
                f"Router edge: {router_edge_pct:.2f}% | harmonic win-rate: {harmonic_win_rate_pct:.2f}%.",
            ],
            "script": (
                f"Our benchmark line is measurable: {benchmark_gain_pct:.2f}% modeled prevention. "
                f"Top lane is {top_sector} at {short_money(top_sector_hourly)} per hour pressure, "
                f"with router edge {router_edge_pct:.2f}% and harmonic win-rate {harmonic_win_rate_pct:.2f}%."
            ),
        },
        {
            "start": "01:00",
            "end": "01:30",
            "title": "Master Valuation",
            "bullets": [
                f"Master valuation proxy: {short_money(valuation_proxy)}.",
                f"Incremental valuation layer: {short_money(valuation_increment)}.",
                f"0.01% math now: {short_money(point_01pct_value)} per year from current measured surface.",
            ],
            "script": (
                f"The current master valuation proxy is {money(valuation_proxy)} with an incremental layer of {money(valuation_increment)}. "
                f"And the scaling math is simple: 0.01% of our currently measured annual value surface is {money(point_01pct_value)} each year."
            ),
        },
        {
            "start": "01:30",
            "end": "02:05",
            "title": "What Investors Are Buying",
            "bullets": [
                f"Grant finder/ranker licensing value: {short_money(grant_license)}.",
                f"Grant/opportunity pipeline value: {short_money(grant_pipeline)}.",
                f"Digital Scout value: {short_money(digital_scout)}.",
                f"Institutional trading value: {short_money(trading_value)}.",
                f"Validated autonomous engine value: {short_money(autonomy_value)}.",
            ],
            "script": (
                "Investors are buying a modular revenue engine. "
                f"Grant platform licensing {money(grant_license)}. Grant pipeline {money(grant_pipeline)}. "
                f"Digital Scout {money(digital_scout)}. Institutional trading system {money(trading_value)}. "
                f"Validated autonomy layer {money(autonomy_value)}."
            ),
        },
        {
            "start": "02:05",
            "end": "02:35",
            "title": "Commercial Packaging",
            "bullets": [
                f"Recommended scenario: {scenario_label}.",
                f"Recommended ARR: {short_money(scenario_arr)}.",
                f"Recommended Year 1 revenue: {short_money(scenario_year1)}.",
                f"Implied EV range: {short_money(scenario_ev_low)} to {short_money(scenario_ev_high)}.",
            ],
            "script": (
                f"Our current recommended commercialization path is {scenario_label}, "
                f"with ARR at {money(scenario_arr)} and Year 1 at {money(scenario_year1)}. "
                f"That maps to an implied enterprise value range of {money(scenario_ev_low)} to {money(scenario_ev_high)}."
            ),
        },
        {
            "start": "02:35",
            "end": "03:00",
            "title": "Close: The Switch",
            "bullets": [
                "This is a machine-verifiable stack with chain-of-custody evidence.",
                "You can fund proof-to-scale, not hype-to-hope.",
                "Founder statement: We are building measurable systems that reduce pain and increase resilience.",
            ],
            "script": (
                "This is the switch. A machine-verifiable system with audited evidence. "
                "You are funding proof-to-scale, not hype-to-hope. "
                "The mission is measurable: reduce pain, reduce failure, and compound value responsibly."
            ),
        },
    ]

    if isinstance(pitch_segments, list) and pitch_segments:
        # Keep continuity with existing mission-control language by appending two strongest legacy lines.
        for row in pitch_segments[:2]:
            if not isinstance(row, dict):
                continue
            line = str(row.get("script") or "").strip()
            if line:
                claim_lines.append(f"Legacy validated pitch line: {line}")

    payload = {
        "generated_utc": now_iso(),
        "schema": "premium_3min_dropmic_deck_v1",
        "title": "Luma Premium Drop-Mic Founder Deck",
        "subtitle": "Founder of the Year Edition | 3-Minute Investor Strike",
        "metrics": {
            "annual_value_signal_usd": annual_value,
            "valuation_proxy_usd": valuation_proxy,
            "valuation_increment_usd": valuation_increment,
            "benchmark_prevented_pct": benchmark_gain_pct,
            "router_edge_pct": router_edge_pct,
            "harmonic_win_rate_pct": harmonic_win_rate_pct,
            "measured_sources": measured_sources,
            "enabled_sources": enabled_sources,
            "measured_coverage_pct": coverage_pct,
            "point_01pct_value_usd": point_01pct_value,
            "point_01pct_twenty_lane_1b_surface_usd": point_01pct_twenty_lanes,
            "top_sector": top_sector,
            "top_sector_hourly_value_usd": top_sector_hourly,
        },
        "valuation_components": {
            "grant_finding_and_ranking_system_license_value_usd": grant_license,
            "grant_and_opportunity_pipeline_value_usd": grant_pipeline,
            "digital_scout_value_usd": digital_scout,
            "institutional_trading_system_value_usd": trading_value,
            "validated_engine_autonomy_value_usd": autonomy_value,
        },
        "commercialization": {
            "recommended_scenario_label": scenario_label,
            "recommended_arr_usd": scenario_arr,
            "recommended_year1_revenue_usd": scenario_year1,
            "recommended_implied_ev_low_usd": scenario_ev_low,
            "recommended_implied_ev_high_usd": scenario_ev_high,
        },
        "patent_anchor": patent_anchor,
        "honest_claim_lines": claim_lines,
        "segments": segments,
    }
    return payload


def apply_holographic_theme(slide) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(8, 16, 36)

    accents = [
        (0.1, 0.2, 2.3, 2.3, RGBColor(31, 91, 255)),
        (7.2, 0.0, 3.0, 3.0, RGBColor(0, 220, 255)),
        (10.4, 5.0, 2.2, 2.2, RGBColor(255, 102, 179)),
    ]
    for left, top, width, height, color in accents:
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 0.72
        shp.line.fill.background()


def add_slide(prs, segment: dict[str, Any], payload: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_holographic_theme(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{segment.get('start')} - {segment.get('end')}  |  {segment.get('title')}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(242, 248, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.8), Inches(0.45))
    stf = subtitle_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = str(payload.get("subtitle", ""))
    sp.font.name = "Segoe UI"
    sp.font.size = Pt(14)
    sp.font.color.rgb = RGBColor(134, 225, 255)

    bullet_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.7),
        Inches(7.9),
        Inches(4.85),
    )
    bullet_box.fill.solid()
    bullet_box.fill.fore_color.rgb = RGBColor(18, 29, 63)
    bullet_box.fill.transparency = 0.10
    bullet_box.line.color.rgb = RGBColor(78, 137, 255)
    bullet_box.line.width = Pt(1.4)

    btf = bullet_box.text_frame
    btf.clear()
    btf.word_wrap = True

    bullets = segment.get("bullets", [])
    if not isinstance(bullets, list):
        bullets = []

    for i, text in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = str(text)
        para.level = 0
        para.space_after = Pt(10)
        para.font.name = "Segoe UI"
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(229, 241, 255)

    metric = payload.get("metrics", {}) if isinstance(payload, dict) else {}
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.8),
        Inches(1.7),
        Inches(3.5),
        Inches(4.85),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(14, 23, 48)
    card.fill.transparency = 0.08
    card.line.color.rgb = RGBColor(0, 223, 255)
    card.line.width = Pt(1.8)

    ctf = card.text_frame
    ctf.clear()
    ctf.word_wrap = True
    lines = [
        f"Value Surface\n{short_money(metric.get('annual_value_signal_usd'))}",
        f"Master Proxy\n{short_money(metric.get('valuation_proxy_usd'))}",
        f"0.01% Math\n{short_money(metric.get('point_01pct_value_usd'))}",
    ]
    for i, text in enumerate(lines):
        para = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        para.text = text
        para.alignment = PP_ALIGN.LEFT
        para.font.name = "Segoe UI"
        para.font.size = Pt(16)
        para.space_after = Pt(10)
        para.font.bold = True if i == 0 else False
        para.font.color.rgb = RGBColor(156, 237, 255)

    footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.8), Inches(11.8), Inches(0.35))
    ftf = footer.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fp.text = "Measured. Auditable. Repeatable. Premium Drop-Mic Founder Deck."
    fp.font.name = "Segoe UI"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(144, 192, 255)

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = str(segment.get("script", "")).strip()


def render_script_markdown(payload: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append("# Premium 3-Minute Drop-Mic Founder Script")
    lines.append("")
    lines.append(f"Generated UTC: {payload.get('generated_utc', '')}")
    lines.append(f"Deck: {payload.get('title', '')}")
    lines.append("")
    lines.append("## Honest Live Numbers")
    lines.append("")
    for row in payload.get("honest_claim_lines", []):
        lines.append(f"- {row}")
    lines.append("")
    lines.append("## Time-Coded Script")
    lines.append("")

    for segment in payload.get("segments", []):
        if not isinstance(segment, dict):
            continue
        lines.append(f"### {segment.get('start', '')} - {segment.get('end', '')} | {segment.get('title', '')}")
        lines.append("")
        lines.append(str(segment.get("script", "")).strip())
        lines.append("")

    lines.append("## Close")
    lines.append("")
    lines.append("Founder line: I am the switch, and this stack is the measurable engine behind that switch.")
    lines.append("Mission line: Build systems that reduce pain, reduce failure, and increase human resilience with measurable proof.")
    lines.append("")
    return "\n".join(lines)


def render_holo_html(payload: dict[str, Any]) -> str:
    metric = payload.get("metrics", {}) if isinstance(payload, dict) else {}
    valuation = payload.get("valuation_components", {}) if isinstance(payload, dict) else {}
    scenario = payload.get("commercialization", {}) if isinstance(payload, dict) else {}
    cards = [
        ("Master Valuation Proxy", short_money(metric.get("valuation_proxy_usd"))),
        ("Valuation Increment", short_money(metric.get("valuation_increment_usd"))),
        ("0.01% of Current Surface", short_money(metric.get("point_01pct_value_usd"))),
        ("Benchmark Prevention", f"{safe_float(metric.get('benchmark_prevented_pct')):.2f}%"),
        ("Grant Platform Licensing", short_money(valuation.get("grant_finding_and_ranking_system_license_value_usd"))),
        ("Recommended ARR", short_money(scenario.get("recommended_arr_usd"))),
    ]

    card_html = "".join(
        f"<div class='card'><div class='k'>{k}</div><div class='v'>{v}</div></div>" for k, v in cards
    )

    segment_html = "".join(
        f"<li><span>{s.get('start','')} - {s.get('end','')}</span> {s.get('title','')}</li>"
        for s in (payload.get("segments", []) if isinstance(payload, dict) else [])
        if isinstance(s, dict)
    )

    return f"""<!doctype html>
<html>
<head>
<meta charset='utf-8' />
<meta name='viewport' content='width=device-width, initial-scale=1' />
<title>Premium Drop-Mic Deck Visuals</title>
<style>
:root {{ --bg:#060c1e; --panel:#102040; --edge:#3ea8ff; --txt:#e9f6ff; --muted:#88bce5; --pink:#f455c4; --cyan:#36f0ff; }}
body {{ margin:0; font-family:Segoe UI, Arial, sans-serif; color:var(--txt); background:radial-gradient(circle at 12% 8%, #1f3f86 0%, transparent 30%), radial-gradient(circle at 88% 20%, #2d0f4a 0%, transparent 32%), var(--bg); }}
.wrap {{ max-width:1200px; margin:0 auto; padding:28px; }}
h1 {{ margin:0 0 8px; font-size:42px; letter-spacing:.02em; }}
.sub {{ color:var(--muted); margin-bottom:18px; }}
.grid {{ display:grid; grid-template-columns:repeat(3,1fr); gap:14px; margin-bottom:20px; }}
.card {{ background:linear-gradient(145deg, rgba(16,32,64,.95), rgba(12,24,48,.86)); border:1px solid rgba(62,168,255,.65); border-radius:16px; padding:14px; box-shadow:0 0 24px rgba(54,240,255,.17); transform:perspective(900px) rotateX(5deg) translateZ(0); }}
.k {{ font-size:11px; letter-spacing:.09em; color:var(--muted); text-transform:uppercase; }}
.v {{ font-size:28px; margin-top:6px; color:var(--cyan); font-weight:700; }}
.panel {{ background:rgba(10,18,38,.82); border:1px solid rgba(244,85,196,.4); border-radius:16px; padding:16px; }}
.panel h2 {{ margin:0 0 10px; font-size:22px; }}
li {{ margin-bottom:8px; }}
li span {{ color:var(--pink); font-weight:700; margin-right:8px; }}
</style>
</head>
<body>
<div class='wrap'>
  <h1>{payload.get('title','Premium Drop-Mic Deck')}</h1>
  <div class='sub'>{payload.get('subtitle','')}</div>
  <div class='grid'>{card_html}</div>
  <div class='panel'>
    <h2>3-Minute Flight Path</h2>
    <ul>{segment_html}</ul>
  </div>
</div>
</body>
</html>
"""


def write_outputs(payload: dict[str, Any]) -> dict[str, str]:
    MISSION_DIR.mkdir(parents=True, exist_ok=True)
    tag = utc_tag()

    pptx_ts = MISSION_DIR / f"premium_dropmic_3min_{tag}.pptx"
    script_ts = MISSION_DIR / f"premium_dropmic_3min_script_{tag}.md"
    json_ts = MISSION_DIR / f"premium_dropmic_3min_pack_{tag}.json"
    html_ts = MISSION_DIR / f"premium_dropmic_3min_visuals_{tag}.html"

    pptx_latest = MISSION_DIR / "premium_dropmic_3min_latest.pptx"
    pptx_status_latest = MISSION_DIR / "premium_dropmic_3min_pptx_status_latest.txt"
    script_latest = MISSION_DIR / "premium_dropmic_3min_script_latest.md"
    json_latest = MISSION_DIR / "premium_dropmic_3min_pack_latest.json"
    html_latest = MISSION_DIR / "premium_dropmic_3min_visuals_latest.html"

    pptx_runtime = ensure_pptx_runtime()
    pptx_latest_out = ""
    if bool(pptx_runtime.get("available", False)):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for segment in payload.get("segments", []):
            if isinstance(segment, dict):
                add_slide(prs, segment, payload)

        prs.save(str(pptx_ts))
        prs.save(str(pptx_latest))
        pptx_latest_out = str(pptx_latest)
        status_line = f"status=ok source={pptx_runtime.get('source','installed')}"
        pptx_status_latest.write_text(status_line, encoding="utf-8")
    else:
        status_line = f"status=skipped reason={pptx_runtime.get('reason','pptx unavailable')}"
        pptx_status_latest.write_text(status_line, encoding="utf-8")

    script_text = render_script_markdown(payload)
    script_ts.write_text(script_text, encoding="utf-8")
    script_latest.write_text(script_text, encoding="utf-8")

    json_text = json.dumps(payload, indent=2)
    json_ts.write_text(json_text, encoding="utf-8")
    json_latest.write_text(json_text, encoding="utf-8")

    html = render_holo_html(payload)
    html_ts.write_text(html, encoding="utf-8")
    html_latest.write_text(html, encoding="utf-8")

    return {
        "pptx_latest": pptx_latest_out,
        "pptx_status": str(pptx_status_latest),
        "script_latest": str(script_latest),
        "json_latest": str(json_latest),
        "visuals_html_latest": str(html_latest),
    }


def main() -> int:
    payload = build_payload()
    paths = write_outputs(payload)
    print(f"PREMIUM_DROPMIC_PPTX={paths['pptx_latest']}")
    print(f"PREMIUM_DROPMIC_PPTX_STATUS={paths['pptx_status']}")
    print(f"PREMIUM_DROPMIC_SCRIPT={paths['script_latest']}")
    print(f"PREMIUM_DROPMIC_PACK_JSON={paths['json_latest']}")
    print(f"PREMIUM_DROPMIC_VISUALS_HTML={paths['visuals_html_latest']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
