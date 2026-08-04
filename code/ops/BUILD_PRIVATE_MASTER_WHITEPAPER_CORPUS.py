"""Build a private, source-traceable corpus for the LumenCore master paper.

The extractor reads local PDF, DOCX, PPTX, Markdown, JSON, and plain-text
sources without modifying them. Patent packages are inventoried but never
expanded. All generated text remains under the ignored ``out`` tree.
"""

from __future__ import annotations

import argparse
import fnmatch
import hashlib
import json
import re
import zipfile
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

from docx import Document
from pypdf import PdfReader
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
ICLOUD = Path.home() / "iCloudDrive"
OUTPUT_ROOT = ROOT / "out" / "private_review" / "master_whitepaper"
EXTRACTED_ROOT = OUTPUT_ROOT / "extracted"
INVENTORY_PATH = OUTPUT_ROOT / "corpus_inventory.json"

STATUS = "PRIVATE_REVIEW_NOT_FOR_PUBLICATION"


@dataclass(frozen=True)
class SourceSpec:
    source_id: str
    base: str
    pattern: str
    intended_role: str
    sensitivity: str
    extract_text: bool = True


SOURCES = (
    SourceSpec(
        "current_source_native_benchmark",
        "repo",
        "docs/LUMENCORE_SOURCE_NATIVE_BENCHMARK_WHITEPAPER_CURRENT.md",
        "authoritative_current_evidence",
        "public_safe_review_required",
    ),
    SourceSpec(
        "master_context_registry",
        "repo",
        "config/luma_master_context_registry_v1.json",
        "architecture_and_lineage",
        "internal",
    ),
    SourceSpec(
        "whitehole_whiteholelab_audit",
        "repo",
        "docs/WHITEHOLE_WHITEHOLELAB_AUDIT_2026-08-02.md",
        "historical_research_provenance_and_claim_boundary",
        "public_safe_review_required",
    ),
    SourceSpec(
        "pitch_deck_governance",
        "repo",
        "docs/PITCH_DECK_GOVERNANCE_2026-07-29.md",
        "presentation_lineage_and_release_control",
        "public_safe_review_required",
    ),
    SourceSpec(
        "nature_inspired_climate_energy",
        "icloud",
        "LumenCore_ A Nature-Inspired, AI-Driven Platform for Climate and Energy Optimization.pdf",
        "concept_history_and_application_map",
        "internal_claim_review",
    ),
    SourceSpec(
        "harmonic_backprop_onepager",
        "icloud",
        "LumenCore_Harmonic-Backprop_OnePager.pdf",
        "algorithm_concept_history",
        "internal_claim_review",
    ),
    SourceSpec(
        "nuclear_harmonization",
        "icloud",
        "LumenCore_Nuclear_Harmonization_WhitePaper.docx",
        "nuclear_application_concept",
        "internal_claim_review",
    ),
    SourceSpec(
        "datacenter_whitepaper",
        "icloud",
        "LumanCore_DataCenter_WhitePaper.pdf",
        "datacenter_application_concept",
        "internal_claim_review",
    ),
    SourceSpec(
        "lumenlogic_fresh_whitepaper",
        "icloud",
        "LumenLogic_Fresh_WhitePaper.docx",
        "legacy_system_thesis",
        "speculative_legacy",
    ),
    SourceSpec(
        "flowform_complete_specs",
        "icloud",
        "Download FLOWFORM Technology Specifications*Complete With Images",
        "hardware_concept_specification",
        "patent_sensitive_review",
    ),
    SourceSpec(
        "flowform_updated_specs",
        "icloud",
        "Download FLOWFORM Technology Specifications*Updated Version",
        "hardware_concept_specification",
        "patent_sensitive_review",
    ),
    SourceSpec(
        "flowform_pdf_specs",
        "icloud",
        "Download FLOWFORM_Technical_Specifications.pdf",
        "hardware_concept_specification",
        "patent_sensitive_review",
    ),
    SourceSpec(
        "lumenshell_flowform_bundle",
        "icloud",
        "LumenShell_FlowForm_Technical_Bundle.pdf",
        "hardware_system_concept",
        "patent_sensitive_review",
    ),
    SourceSpec(
        "lumenshell_spec_01",
        "icloud",
        "LumenShell-SPEC-01.pdf",
        "hardware_system_concept",
        "patent_sensitive_review",
    ),
    SourceSpec(
        "aetherframe_genesis_logbook",
        "icloud",
        "repaired_Download AetherFrame_Genesis_Logbook.pdf",
        "concept_history",
        "internal_claim_review",
    ),
    SourceSpec(
        "aetherframe_pitch_deck",
        "icloud",
        "Download AetherFrame Pitch Deck 2",
        "application_and_product_history",
        "internal_claim_review",
    ),
    SourceSpec(
        "novacore_concept_paper",
        "icloud",
        "NovaCore_Concept_Paper.pdf",
        "concept_history",
        "speculative_legacy",
    ),
    SourceSpec(
        "agentic_ai_blueprint",
        "icloud",
        "LumenKing_Agentic_AI_Blueprint.pdf",
        "agent_architecture_concept",
        "internal_claim_review",
    ),
    SourceSpec(
        "telecom_deck",
        "icloud",
        "LumenLogic_Telecom_Deck.pptx",
        "telecom_application_concept",
        "internal_claim_review",
    ),
    SourceSpec(
        "wormhole_field_research",
        "icloud",
        "LumenCore_Wormhole_Field_Research_WhitePaper 2.pdf",
        "historical_speculative_appendix",
        "speculative_legacy",
    ),
    SourceSpec(
        "wormhole_grant_proposal",
        "icloud",
        "LumenCore_Wormhole_Research_Proposal_Gov_Grant_Format.pdf",
        "historical_speculative_appendix",
        "speculative_legacy",
    ),
    SourceSpec(
        "logic_specs_placeholder",
        "icloud",
        "LumenCore_LogicSpecs.docx",
        "placeholder_audit",
        "internal",
    ),
    SourceSpec(
        "lumencore_patent_bundle",
        "icloud",
        "LumenCore_Patent_Bundle.zip",
        "existence_and_lineage_only",
        "patent_package_do_not_expand",
        False,
    ),
    SourceSpec(
        "lumenshell_patent_package",
        "icloud",
        "LumenShell_Patent_Package.zip",
        "existence_and_lineage_only",
        "patent_package_do_not_expand",
        False,
    ),
)


FLAG_PATTERNS = {
    "contact_identifier": re.compile(
        r"(?:[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}|"
        r"(?:\+?1[-. ]?)?\(?\d{3}\)?[-. ]\d{3}[-. ]\d{4})",
        re.IGNORECASE,
    ),
    "credential_language": re.compile(
        r"\b(password|passcode|one[- ]time code|otp|access token|api key|secret key)\b",
        re.IGNORECASE,
    ),
    "government_or_legal_identifier": re.compile(
        r"\b(UEI|CAGE|EIN|application\s*(?:number|#)|filing\s*(?:number|#))\b",
        re.IGNORECASE,
    ),
    "street_address": re.compile(
        r"\b\d{2,5}\s+(?:(?-i:[A-Z])[A-Z0-9.'-]{0,30}\s+){1,6}"
        r"(?:ST|STREET|AVE|AVENUE|RD|ROAD|DR|DRIVE|LN|LANE|BLVD|BOULEVARD)\b",
        re.IGNORECASE,
    ),
    "patent_language": re.compile(
        r"\b(patent|provisional|claims?|inventor|prior art|trade secret|proprietary)\b",
        re.IGNORECASE,
    ),
    "speculative_physics": re.compile(
        r"\b(wormhole|zero[- ]point|scalar field|weather control|"
        r"consciousness field|free energy|anti[- ]gravity)\b",
        re.IGNORECASE,
    ),
    "unverified_benefit_language": re.compile(
        r"\b(guaranteed|revolutionary|breakthrough|world[- ]first|"
        r"trillion[- ]dollar|percent savings|eliminates? risk)\b",
        re.IGNORECASE,
    ),
    "authorship_or_boilerplate_conflict": re.compile(
        r"\b(sammuti\.com|need professional help in developing your architecture)\b",
        re.IGNORECASE,
    ),
}


def sha256_bytes(payload: bytes) -> str:
    return hashlib.sha256(payload).hexdigest()


def safe_source_path(spec: SourceSpec) -> Path | None:
    base = ROOT if spec.base == "repo" else ICLOUD
    if not any(character in spec.pattern for character in "*?["):
        candidate = base / spec.pattern
        return candidate if candidate.is_file() else None
    matches = [
        path
        for path in base.iterdir()
        if path.is_file() and fnmatch.fnmatchcase(path.name, spec.pattern)
    ]
    return sorted(matches, key=lambda path: path.name.casefold())[0] if matches else None


def detect_kind(path: Path, payload: bytes) -> str:
    suffix = path.suffix.casefold()
    if payload.startswith(b"%PDF-"):
        return "pdf"
    if payload.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(path) as archive:
                names = set(archive.namelist())
        except (OSError, zipfile.BadZipFile):
            return "invalid_zip"
        if "word/document.xml" in names:
            return "docx"
        if "ppt/presentation.xml" in names:
            return "pptx"
        return "zip"
    if suffix in {".md", ".txt", ".json"}:
        return suffix.lstrip(".")
    return "unknown"


def iter_docx_text(path: Path) -> Iterable[str]:
    document = Document(path)
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            yield paragraph.text
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                yield " | ".join(values)


def iter_pptx_text(path: Path) -> Iterable[str]:
    presentation = Presentation(path)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        yield f"[Slide {slide_index}]"
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                yield text.strip()


def extract_text(path: Path, kind: str) -> tuple[str, int | None]:
    if kind == "pdf":
        reader = PdfReader(path)
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(
            f"[Page {index}]\n{page_text.strip()}"
            for index, page_text in enumerate(pages, start=1)
        )
        return text.strip(), len(pages)
    if kind == "docx":
        return "\n\n".join(iter_docx_text(path)).strip(), None
    if kind == "pptx":
        presentation = Presentation(path)
        return "\n\n".join(iter_pptx_text(path)).strip(), len(presentation.slides)
    if kind in {"md", "txt", "json"}:
        return path.read_text(encoding="utf-8", errors="replace").strip(), None
    return "", None


def flag_text(text: str) -> list[str]:
    return [name for name, pattern in FLAG_PATTERNS.items() if pattern.search(text)]


def build_inventory() -> dict[str, object]:
    EXTRACTED_ROOT.mkdir(parents=True, exist_ok=True)
    records: list[dict[str, object]] = []
    for spec in SOURCES:
        path = safe_source_path(spec)
        if path is None:
            records.append(
                {
                    **asdict(spec),
                    "status": "MISSING",
                    "resolved_path": None,
                }
            )
            continue

        payload = path.read_bytes()
        record: dict[str, object] = {
            **asdict(spec),
            "status": "PRESENT",
            "resolved_path": str(path),
            "bytes": len(payload),
            "source_sha256": sha256_bytes(payload),
            "modified_utc": datetime.fromtimestamp(
                path.stat().st_mtime, tz=timezone.utc
            ).isoformat(),
        }

        if not spec.extract_text:
            record["status"] = "INVENTORIED_NOT_EXPANDED"
            record["kind"] = "opaque_archive"
            record["archive_contents_inspected"] = False
            record["archive_member_count"] = None
            records.append(record)
            continue

        kind = detect_kind(path, payload)
        record["kind"] = kind

        if kind not in {"pdf", "docx", "pptx", "md", "txt", "json"}:
            record["status"] = "UNREADABLE_OR_PLACEHOLDER"
            records.append(record)
            continue

        try:
            text, page_or_slide_count = extract_text(path, kind)
        except Exception as exc:  # document parsers raise several format errors
            record["status"] = "EXTRACTION_FAILED"
            record["error_type"] = type(exc).__name__
            records.append(record)
            continue

        record["extracted_characters"] = len(text)
        record["extracted_words"] = len(text.split())
        record["page_or_slide_count"] = page_or_slide_count
        record["content_flags"] = flag_text(text)
        if len(text) < 80:
            record["status"] = "EMPTY_OR_PLACEHOLDER_TEXT"
        extracted_path = EXTRACTED_ROOT / f"{spec.source_id}.txt"
        extracted_path.write_text(
            f"STATUS: {STATUS}\nSOURCE: {path}\nSHA256: {record['source_sha256']}\n\n{text}\n",
            encoding="utf-8",
        )
        record["extracted_path"] = str(extracted_path)
        record["extracted_sha256"] = sha256_bytes(extracted_path.read_bytes())
        records.append(record)

    present_count = sum(
        1 for record in records if record.get("status") not in {"MISSING"}
    )
    extractable_count = sum(
        1
        for record in records
        if isinstance(record.get("extracted_characters"), int)
        and int(record["extracted_characters"]) >= 80
    )
    inventory: dict[str, object] = {
        "schema": "lumencore_private_master_whitepaper_corpus_v1",
        "status": STATUS,
        "generated_utc": datetime.now(timezone.utc).isoformat(),
        "source_count": len(records),
        "present_count": present_count,
        "extractable_count": extractable_count,
        "publication_authorized": False,
        "patent_archives_expanded": False,
        "records": records,
    }
    inventory["inventory_sha256"] = sha256_bytes(
        json.dumps(
            inventory,
            sort_keys=True,
            ensure_ascii=True,
            separators=(",", ":"),
        ).encode("utf-8")
    )
    return inventory


def verify_inventory(inventory: dict[str, object]) -> None:
    if inventory.get("schema") != "lumencore_private_master_whitepaper_corpus_v1":
        raise ValueError("Unexpected corpus inventory schema")
    if inventory.get("status") != STATUS:
        raise ValueError("Corpus inventory is not private-review gated")
    if inventory.get("publication_authorized") is not False:
        raise ValueError("Corpus inventory publication gate must remain false")
    if inventory.get("patent_archives_expanded") is not False:
        raise ValueError("Patent archive expansion must remain false")
    expected_hash = inventory.get("inventory_sha256")
    unsigned = dict(inventory)
    unsigned.pop("inventory_sha256", None)
    actual_hash = sha256_bytes(
        json.dumps(
            unsigned,
            sort_keys=True,
            ensure_ascii=True,
            separators=(",", ":"),
        ).encode("utf-8")
    )
    if expected_hash != actual_hash:
        raise ValueError("Corpus inventory hash mismatch")
    for record in inventory.get("records", []):
        if not isinstance(record, dict) or record.get("status") == "MISSING":
            continue
        path_value = record.get("resolved_path")
        if not isinstance(path_value, str):
            raise ValueError("Present source has no resolved path")
        path = Path(path_value)
        if not path.is_file() or sha256_bytes(path.read_bytes()) != record.get(
            "source_sha256"
        ):
            raise ValueError(f"Source receipt mismatch: {path}")
        if record.get("extract_text") is False:
            if record.get("status") != "INVENTORIED_NOT_EXPANDED":
                raise ValueError(f"Opaque archive has invalid status: {path}")
            if record.get("kind") != "opaque_archive":
                raise ValueError(f"Opaque archive kind drifted: {path}")
            if record.get("archive_contents_inspected") is not False:
                raise ValueError(f"Opaque archive inspection must remain false: {path}")
            if record.get("archive_member_count") is not None:
                raise ValueError(f"Opaque archive member count must remain absent: {path}")
            forbidden = {
                "extracted_path",
                "extracted_sha256",
                "extracted_characters",
                "extracted_words",
                "page_or_slide_count",
                "content_flags",
            }
            if forbidden.intersection(record):
                raise ValueError(f"Opaque archive leaked extracted metadata: {path}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--check", action="store_true")
    args = parser.parse_args()
    if args.check:
        inventory = json.loads(INVENTORY_PATH.read_text(encoding="utf-8"))
        verify_inventory(inventory)
    else:
        inventory = build_inventory()
        OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
        INVENTORY_PATH.write_text(
            json.dumps(inventory, indent=2, ensure_ascii=True) + "\n",
            encoding="utf-8",
        )
        verify_inventory(inventory)
    print(
        json.dumps(
            {
                "schema": inventory["schema"],
                "status": inventory["status"],
                "source_count": inventory["source_count"],
                "present_count": inventory["present_count"],
                "extractable_count": inventory["extractable_count"],
                "publication_authorized": inventory["publication_authorized"],
                "inventory_path": str(INVENTORY_PATH),
            },
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
